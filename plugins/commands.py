# Don't Remove Credit Tg - @VJ_Botz
# Subscribe YouTube Channel For Amazing Bot https://youtube.com/@Tech_VJ
# Ask Doubt on telegram @KingVJ01

import os
import logging
import random
import asyncio
from validators import domain
from Script import script
from plugins.dbusers import db
from pyrogram import Client, filters, enums
from plugins.users_api import get_user, update_user_info
from pyrogram.errors import ChatAdminRequired, FloodWait
from pyrogram.types import *
from utils import verify_user, check_token, check_verification, get_token
from config import *
import re
import json
import base64
from urllib.parse import quote_plus
import os
from PIL import Image
from fpdf import FPDF
import img2pdf
from io import BytesIO
from PyPDF2 import PdfReader, PageObject
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from docx import Document
from fpdf import FPDF
from pptx import Presentation
from datetime import datetime, timedelta, date
import pytz
from pypdf import PdfReader, PdfWriter
from TechVJ.utils.file_properties import get_name, get_hash, get_media_file_size
import os
from io import BytesIO
import img2pdf
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from pdf2image import convert_from_path
from reportlab.pdfgen import canvas
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
from reportlab.lib.pagesizes import letter

from pyrogram import Client, filters
from pyrogram.types import InlineKeyboardMarkup, InlineKeyboardButton, CallbackQuery

logger = logging.getLogger(__name__)

BATCH_FILES = {}
user_images = {}
user_docs = {}
user_pdfs = {}
last_message = {}

# Supported formats
IMAGE_FORMATS = (".jpg", ".jpeg", ".png")
PDF_FORMATS = (".pdf",)
DOC_FORMATS = (".txt", ".docx", ".pptx")


# Don't Remove Credit Tg - @VJ_Botz
# Subscribe YouTube Channel For Amazing Bot https://youtube.com/@Tech_VJ
# Ask Doubt on telegram @KingVJ01

@Client.on_message(filters.command("user_details"))
async def user_details(client, message):
    try:
        if len(message.command) > 1:
            # Command with user_id argument: /user_details 12345
            user_id = int(message.command[1])
        else:
            # Command without user_id argument: /user_details
            user_id = message.from_user.id
            
        user_data = await db.ud.find_one({"id": user_id})
        
        if user_data:
            details_message = (
                f"User Details:\n\n"
                f"Joined on: {user_data.get('joined', 'N/A')}\n"
                f"ID: {user_data.get('id', 'N/A')}\n"
                f"Img2PDF: {user_data.get('I2P', 0)}\n"
                f"PDF Watermark: {user_data.get('PW', 0)}\n"
                f"PDF2IMG: {user_data.get('P2I', 0)}\n"
                f"PDF Protect: {user_data.get('PP', 0)}\n"
                f"PDF Merge: {user_data.get('PM', 0)}"
            )
            await message.reply(details_message)
        else:
            await message.reply("User details not found.")
            
    except ValueError:
        await message.reply("Invalid user ID. Please provide a valid numerical user ID.")
    except Exception as e:
        await message.reply(f"An error occurred: {str(e)}")
def get_size(size):
    """Get size in readable format"""

    units = ["Bytes", "KB", "MB", "GB", "TB", "PB", "EB"]
    size = float(size)
    i = 0
    while size >= 1024.0 and i < len(units):
        i += 1
        size /= 1024.0
    return "%.2f %s" % (size, units[i])

def formate_file_name(file_name):
    chars = ["[", "]", "(", ")"]
    for c in chars:
        file_name.replace(c, "")
    file_name = '@VJ_Botz ' + ' '.join(filter(lambda x: not x.startswith('http') and not x.startswith('@') and not x.startswith('www.'), file_name.split()))
    return file_name

# Don't Remove Credit Tg - @VJ_Botz
# Subscribe YouTube Channel For Amazing Bot https://youtube.com/@Tech_VJ
# Ask Doubt on telegram @KingVJ0

@Client.on_message(filters.command("clear") & filters.incoming)
async def clear(client, message):
    try:
        for file_path in user_pdfs.get(message.from_user.id, []):
            if os.path.exists(file_path):
                os.remove(file_path)
        user_pdfs[message.from_user.id] = []
    except Exception as e:
        await message.reply_text(f"An error occurred: {e}")



@Client.on_message(filters.command("start") & filters.incoming)
async def start(client, message):
    try:
        user_id = message.from_user.id 
        if not await db.is_user_exist(message.from_user.id):
            await db.add_user(message.from_user.id, message.from_user.first_name)
            await client.send_message(LOG_CHANNEL, script.LOG_TEXT.format(message.from_user.id, message.from_user.mention))
   
        if not await db.ud.find_one({"id": user_id}):
            user_data = {
                "id": user_id,
                "joined": await dati()
            }    
            await db.ud.update_one({"id": user_data["id"]}, {"$set": user_data}, upsert=True)
            
        await message.reply("welcome")
    except Exception as e:
        await message.reply_text(f"An error occurred: {e}")



async def dati():
    try:
        kolkata_timezone = pytz.timezone('Asia/Kolkata')
        kolkata_time = datetime.now(kolkata_timezone)
        formatted_time = kolkata_time.strftime('%d/%m/%Y %H:%M:%S')  
        return formatted_time 
    except Exception as e:
        # Handle exceptions appropriately, e.g., logging or raising
        raise RuntimeError(f"Error in dati function: {str(e)}")

@Client.on_message(filters.command("rate_me"))
async def rate_me(client, message):
    user_id = message.from_user.id
    user_rating = await db.get_user_rating(user_id)
    if user_rating is not None:
        await message.reply_text(f'You have already rated this bot: {user_rating} stars.\nIf you want to re-rate, please use /re_rate_me.')
    else:
        await rating_poll(client, message, user_id, re_rating=False)

@Client.on_message(filters.command("re_rate_me"))
async def re_rate_me(client, message):
    user_id = message.from_user.id
    await rating_poll(client, message, user_id, re_rating=True)



# Initialize storage for images, PDFs, and last messages


### Collect images and PDFs separately ###
@Client.on_message(
    filters.private & 
    (filters.photo | filters.document) & 
    filters.incoming
)
async def collect_files(bot, message):
    try:
        user_id = message.from_user.id

        # Initialize user storage if not present
        if user_id not in user_images:
            user_images[user_id] = []
        if user_id not in user_pdfs:
            user_pdfs[user_id] = []
        if user_id not in user_docs:
            user_docs[user_id] = []

        # Validate file type correctly
        if message.document:
            file_name = message.document.file_name.lower() if message.document.file_name else ""

            # Define valid extensions
            IMAGE_FORMATS = (".jpg", ".jpeg", ".png")
            PDF_FORMATS = (".pdf",)
            DOC_FORMATS = (".txt", ".docx", ".pptx")

            # Categorize the file
            if file_name.endswith(IMAGE_FORMATS):
                file_path = await message.download()
                user_images[user_id].append(file_path)
            elif file_name.endswith(PDF_FORMATS):
                file_path = await message.download()
                user_pdfs[user_id].append(file_path)
            elif file_name.endswith(DOC_FORMATS):
                file_path = await message.download()
                user_docs[user_id].append(file_path)
            else:
                await message.reply_text("Unsupported file type. Please send a valid file.")
                return

        # Delete the last message if exists
        if user_id in last_message and last_message[user_id]:
            try:
                await last_message[user_id][-1].delete()
                last_message[user_id].pop()
            except Exception as e:
                print(f"Failed to delete the last message: {e}")

        # Create dynamic buttons based on uploaded files
        buttons = []
        
        if user_images[user_id]:
            buttons.append([InlineKeyboardButton("Create PDF", callback_data="create_pdf")])

        if user_docs[user_id]:
            for doc in user_docs[user_id]:
                if doc.endswith('.docx') or doc.endswith('.pptx'):
                    buttons.append([InlineKeyboardButton("Convert Documents to PDF", callback_data="convert_docs")])
                elif doc.endswith('.txt'):
                    buttons.append([InlineKeyboardButton("Convert TXT to PDF", callback_data="convert_txt")])

        if len(user_pdfs[user_id]) == 1:
            buttons.append([InlineKeyboardButton("Watermark PDF", callback_data="watermark_pdf")])
            buttons.append([InlineKeyboardButton("Protect PDF", callback_data="prot_pdf")])  # New Button

        # Send reply with buttons
        sent_message = await message.reply_text(
            f"You have {len(user_images[user_id])} image(s) and {len(user_pdfs[user_id])} PDF(s) ready.\n\n"
            "Choose an action to proceed.", 
            reply_markup=InlineKeyboardMarkup(buttons)
        )

        if user_id not in last_message:
            last_message[user_id] = []
        last_message[user_id].append(sent_message)

    except Exception as e:
        await message.reply_text(f"An error occurred: {e}")


### Handle Button Actions ###
@Client.on_callback_query()
async def cb_handler(client: Client, query: CallbackQuery):
    user_id = query.from_user.id

    if query.data == "create_pdf":
        await create_pdf(client, query, user_id)

    elif query.data == "convert_docs":
        await convert_docs_to_pdf(client, query, user_id)
    
    elif query.data == "convert_pptx":
        await convert_pptx_to_pdf(client, query, user_id)

    elif query.data == "convert_txt":
        await convert_txt_to_pdf(client, query, user_id)


    elif query.data == "prot_pdf":
        await protect_pdf(client, query, user_id)

    elif query.data == "extract_images":
        await select_image_format(client, query, user_id)

    elif query.data.startswith("extract_format_"):
        image_format = query.data.split("_")[-1]
        await extract_images(client, query, user_id, image_format)

    
    elif query.data == "watermark_pdf":
        await ask_watermark_options(client, query, user_id)  # Start with watermark type selection

    elif query.data.startswith("watermark_position_"):
        position = query.data.split("_")[-1]
        watermark_data = user_watermark_data.get(user_id)  # Retrieve stored watermark data
        if watermark_data:
            await watermark_pdf(client, query, user_id, position, watermark_data)
        else:
            await query.message.edit_text("Error: Watermark data missing!")

    elif query.data.startswith("watermark_type_"):
        watermark_type = query.data.split("_")[-1]
        await ask_watermark_details(client, query, user_id, watermark_type)

    
    elif query.data == "merge_pdfs":
        await merge_pdfs(client, query, user_id)


### Create PDF from Images ###
async def create_pdf(client, query, user_id):
    await query.message.edit_text("Converting images to PDF, please wait...")
    image_files = user_images.get(user_id, [])
    if not image_files:
        await query.answer("No images available for PDF creation.", show_alert=True)
        return

    try:
        response = await client.ask(user_id, "Send the PDF name (without extension):")
        pdf_name = "".join(char for char in response.text if char.isalnum() or char in " _-").strip()
        pdf_file_name = f"{pdf_name or 'converted'}.pdf"

        pdf_bytes = img2pdf.convert(image_files)
        pdf_output = BytesIO(pdf_bytes)
        pdf_output.seek(0)

        await client.send_document(
            chat_id=user_id,
            document=pdf_output,
            file_name=pdf_file_name,
            caption=f"Here is your PDF: **{pdf_file_name}**"
        )
        
        user_data = await db.ud.find_one({"id": user_id})
        user_data["I2P"] = user_data.get("I2P", 0) + 1
        await db.ud.update_one({"id": user_data["id"]}, {"$set": {"I2P": user_data["I2P"]}}, upsert=True)

        clear_user_data(user_id, "images")

    except Exception as e:
        await query.answer(f"Error while creating PDF: {str(e)}", show_alert=True)


### Select Image Format for Extraction ###
async def select_image_format(client, query, user_id):
    buttons = [
        [InlineKeyboardButton("JPG", callback_data="extract_format_jpg")],
        [InlineKeyboardButton("PNG", callback_data="extract_format_png")]
    ]
    await query.message.edit_text(
        "Select the image format for extraction:",
        reply_markup=InlineKeyboardMarkup(buttons)
    )


### Extract Images from PDF ###
async def extract_images(client, query, user_id, image_format):
    await query.message.edit_text("Extracting images from PDF, please wait...")
    try:
        for pdf_path in user_pdfs.get(user_id, []):
            images = convert_from_path(pdf_path, fmt=image_format, thread_count=4)
            pdf_name = os.path.splitext(os.path.basename(pdf_path))[0]
            for idx, img in enumerate(images):
                img_bytes = BytesIO()
                img.save(img_bytes, format=image_format.upper())
                img_bytes.seek(0)

                await client.send_document(
                    chat_id=user_id,
                    document=img_bytes,
                    file_name=f"{pdf_name}_page_{idx+1}.{image_format}"
                )
                
                user_data = await db.ud.find_one({"id": user_id})
                user_data["P2I"] = user_data.get("P2I", 0) + 1
                await db.ud.update_one({"id": user_data["id"]}, {"$set": {"P2I": user_data["P2I"]}}, upsert=True)
   
                clear_user_data(user_id, "pdfs")
                
    except Exception as e:
        await message.reply_text(f"An error occurred: {e}")


from pyrogram.types import InlineKeyboardMarkup, InlineKeyboardButton
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
import os

# Store watermark data for each user
user_watermark_data = {}



async def ask_watermark_options(client, query, user_id):
    """ Ask users whether they want text, image, or both watermarks """
    buttons = [
        [InlineKeyboardButton("Text Only", callback_data="watermark_type_text")],
        [InlineKeyboardButton("Image Only", callback_data="watermark_type_image")],
        [InlineKeyboardButton("Text & Image", callback_data="watermark_type_both")]
    ]
    await query.message.edit_text(
        "Select the type of watermark you want:",
        reply_markup=InlineKeyboardMarkup(buttons)
    )


async def ask_watermark_details(client, query, user_id, watermark_type):
    """ Ask for text or image based on selection """
    try:
        watermark_data = {"text": None, "image": None}

        if watermark_type in ["text", "both"]:
            response = await client.ask(user_id, "Send watermark text:")
            watermark_data["text"] = response.text or "Poojar Project"

        if watermark_type in ["image", "both"]:
            response = await client.ask(user_id, "Send watermark image:")
            if response.photo:
                photo = response.photo  # Fix: No need to use [-1]
                file = await client.download_media(photo.file_id)  # Fix: Use .file_id
                watermark_data["image"] = file

        # Store watermark data for later use
        user_watermark_data[user_id] = watermark_data

        await select_watermark_position(client, query, user_id)

    except Exception as e:
        await query.message.reply_text(f"An error occurred: {e}")


async def select_watermark_position(client, query, user_id):
    """ Ask users where they want the watermark """
    buttons = [
        [InlineKeyboardButton("Top", callback_data="watermark_position_top")],
        [InlineKeyboardButton("Center Cross", callback_data="watermark_position_center")],
        [InlineKeyboardButton("Bottom", callback_data="watermark_position_bottom")],
        [InlineKeyboardButton("Cross (full page)", callback_data="watermark_position_cross")]
    ]
    await query.message.edit_text(
        "Select where you want the watermark:",
        reply_markup=InlineKeyboardMarkup(buttons)
    )

async def watermark_pdf(client, query, user_id, position, watermark_data):
    """ Process and apply watermark to PDFs """
    await query.message.edit_text("Watermarking PDF, please wait...")
    try:
        text = watermark_data.get("text")
        image_path = watermark_data.get("image")

        for pdf_path in user_pdfs.get(user_id, []):
            reader = PdfReader(pdf_path)
            writer = PdfWriter()

            for page_num, page in enumerate(reader.pages):
                page_width = float(page.mediabox.width)
                page_height = float(page.mediabox.height)

                watermark_buffer = create_watermark_pdf(
                    text, position, page_width, page_height, image_path
                )

                watermark_reader = PdfReader(watermark_buffer)
                watermark_page = watermark_reader.pages[0]

                # Fix for PyPDF2 v3+
                page.merge_page(watermark_page)
                writer.add_page(page)

            output_path = f"{os.path.splitext(pdf_path)[0]}_watermarked.pdf"
            with open(output_path, "wb") as f_out:
                writer.write(f_out)

            await client.send_document(user_id, document=output_path)

    except Exception as e:
        await query.message.edit_text(f"Error during watermarking: {e}")

from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
from PyPDF2 import PdfReader, PdfWriter

def create_watermark_pdf(text, position, page_width, page_height, image_path=None):
    """Create a watermark PDF in memory"""
    try:
        pdf_buffer = BytesIO()
        c = canvas.Canvas(pdf_buffer, pagesize=(page_width, page_height))

        text_font_size = min(page_width, page_height) * 0.05  
        cross_font_size = min(page_width, page_height) * 0.08  
        center_font_size = min(page_width, page_height) * 0.12  

        c.setFont("Helvetica-Bold", text_font_size)
        c.setFillColor(Color(0, 0, 0, alpha=0.3))

        # Adjusted Y positions for visibility
        pos = {
            "top": (page_width / 2, page_height - (text_font_size * 2)),  
            "center": (page_width / 2, page_height / 2),
            "bottom": (page_width / 2, text_font_size * 2),
        }

        x, y = pos.get(position, (page_width / 2, page_height / 2))

        if text:
            c.saveState()
            if position == "center":
                c.setFont("Helvetica-Bold", center_font_size)
                c.translate(x, y)
                c.rotate(45)
                c.drawCentredString(0, 0, text)
            elif position == "cross":
                c.setFont("Helvetica-Bold", cross_font_size)
                text_width = c.stringWidth(text, "Helvetica-Bold", cross_font_size)  # Get actual text width
                step_size_x = text_width + 70  # Adjust spacing (56.7 to 85 for 2-3 cm)
                step_size_y = 100  # Adjust spacing (85 to 142 for 3-5 cm)
                for i in range(-int(page_width // step_size_x), int(page_width // step_size_x) + 2):
                    for j in range(-int(page_height // step_size_y), int(page_height // step_size_y) + 2):
                        c.saveState()
                        cx = i * step_size_x
                        cy = j * step_size_y
                        c.translate(cx, cy)
                        c.rotate(30)
                        c.drawCentredString(0, 0, text)
                        c.restoreState()
            else:
                # Debug print coordinates to check placement
                print(f"Placing text at: {x}, {y} for position {position}")
                c.drawCentredString(x, y, text)

        if image_path:
            image = ImageReader(image_path)
            img_original_width, img_original_height = image.getSize()
            scale_factor = min(page_width * 0.4 / img_original_width, page_height * 0.4 / img_original_height)
            img_width = img_original_width * scale_factor
            img_height = img_original_height * scale_factor
            img_x = (page_width - img_width) / 2
            img_y = (page_height - img_height) / 2
            c.drawImage(image, img_x, img_y, img_width, img_height, mask="auto")

        c.save()
        pdf_buffer.seek(0)
        return pdf_buffer
    except Exception as e:
        print(f"Error in create_watermark_pdf: {e}")
        return None


    
        

import os
from PyPDF2 import PdfReader, PdfWriter

async def protect_pdf(client, query, user_id):
    await query.message.edit_text("Preparing to protect your PDF...")

    if len(user_pdfs[user_id]) == 1:
        try:
            # Ask for the password
            response = await client.ask(user_id, "Send the password to protect the PDF:")
            password = response.text or "Poojar Project"

            input_pdf_path = user_pdfs[user_id][0]
            # Generate the output path properly
            output_pdf_path = os.path.join(
                os.path.dirname(input_pdf_path), 
                f"{os.path.basename(input_pdf_path)}"
            )

            # Apply password protection
            writer = PdfWriter()
            with open(input_pdf_path, "rb") as file:
                reader = PdfReader(file)
                for page_num in range(len(reader.pages)):
                    writer.add_page(reader.pages[page_num])
                
                writer.encrypt(password)
                
                with open(output_pdf_path, "wb") as output_file:
                    writer.write(output_file)
            
            await client.send_document(user_id, output_pdf_path, caption="Your password-protected PDF is ready!")
    
            user_data = await db.ud.find_one({"id": user_id})
            user_data["PP"] = user_data.get("PP", 0) + 1
            await db.ud.update_one({"id": user_data["id"]}, {"$set": {"PP": user_data["PP"]}}, upsert=True)



            clear_user_data(user_id, "pdfs")
            
            os.remove(output_pdf_path)  # Corrected variable name

        except Exception as e:
            await client.send_message(user_id, f"Failed to protect PDF: {e}")


from fpdf import FPDF
import os

async def convert_txt_to_pdf(client, query, user_id):
    await query.message.edit_text("Converting .txt file to PDF, please wait...")

    txt_files = user_docs.get(user_id, [])
    if not txt_files:
        await query.answer("No .txt files available for PDF conversion.", show_alert=True)
        return

    try:
        response = await client.ask(user_id, "Send the PDF name (without extension):")
        pdf_name = "".join(char for char in response.text if char.isalnum() or char in " _-").strip()
        pdf_file_name = f"{pdf_name or 'converted'}.pdf"

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.set_font("Arial", size=12)

        for txt_file in txt_files:
            if txt_file.lower().endswith(".txt"):
                pdf.add_page()

                with open(txt_file, "r", encoding="utf-8", errors="ignore") as file:
                    for line in file:
                        pdf.multi_cell(0, 10, line.strip())

        temp_pdf_path = "{}".format(pdf_file_name)
        pdf.output(temp_pdf_path)

        await client.send_document(
            chat_id=user_id,
            document=temp_pdf_path,
            caption=f"Here is your converted PDF: **{pdf_file_name}**"
        )

        clear_user_data(user_id, "docs")

    except Exception as e:
        await query.answer(f"Error while converting .txt to PDF: {str(e)}", show_alert=True)

import os
from fpdf import FPDF
from pptx import Presentation
from docx import Document
from PIL import Image, ImageDraw
from io import BytesIO
import aiofiles

async def convert_docs_to_pdf(client, query, user_id):
    await query.message.edit_text("Converting documents to PDF, please wait...")

    doc_files = user_docs.get(user_id, [])
    if not doc_files:
        await query.answer("No documents available for PDF conversion.", show_alert=True)
        return

    try:
        response = await client.ask(user_id, "Send the PDF name (without extension):")
        pdf_name = "".join(char for char in response.text if char.isalnum() or char in " _-").strip()
        pdf_file_name = f"{pdf_name or 'converted'}.pdf"

        image_files = []

        for doc_file in doc_files:
            
            ### Convert PPTX to Images ###
            if doc_file.lower().endswith(".pptx"):
                await query.message.edit_text("Processing PPTX file...")
                ppt = Presentation(doc_file)
                for i, slide in enumerate(ppt.slides):
                    slide_image_path = f"slide_{i}.png"
                    slide_image = Image.new("RGB", (1280, 720), "white")
                    draw = ImageDraw.Draw(slide_image)
                    
                    # Extract and render text from slides
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)

                    draw.multiline_text((50, 50), "\n".join(text_content), fill="black")
                    slide_image.save(slide_image_path)
                    image_files.append(slide_image_path)
            
            ### Convert DOCX to Images ###
            elif doc_file.lower().endswith(".docx"):
                await query.message.edit_text("Processing DOCX file...")
                doc = Document(doc_file)
                for i, para in enumerate(doc.paragraphs):
                    page_image_path = f"page_{i}.png"
                    page_image = Image.new("RGB", (595, 842), "white")  # A4 size at 72 DPI
                    draw = ImageDraw.Draw(page_image)
                    
                    # Render the paragraph text properly
                    draw.multiline_text((50, 50), para.text, fill="black")
                    page_image.save(page_image_path)
                    image_files.append(page_image_path)

        ### Convert Images to PDF ###
        if not image_files:
            await query.answer("No images generated for PDF creation.", show_alert=True)
            return

        await query.message.edit_text("Generating PDF file...")

        pdf = FPDF()
        for image_file in image_files:
            pdf.add_page()
            pdf.image(image_file, x=0, y=0, w=210, h=297)  # A4 dimensions in mm

        # Save the PDF to a temporary file asynchronously
        temp_pdf_path = f"/tmp/{pdf_file_name}"
        pdf.output(temp_pdf_path, 'F')

        async with aiofiles.open(temp_pdf_path, 'rb') as pdf_file:
            await client.send_document(
                chat_id=user_id,
                document=pdf_file,
                file_name=pdf_file_name,
                caption=f"Here is your converted PDF: **{pdf_file_name}**"
            )

        clear_user_data(user_id, "docs")

        # Clean up image files and temporary PDF
        for img in image_files:
            if os.path.exists(img):
                os.remove(img)
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

    except Exception as e:
        await query.answer(f"Error while converting documents: {str(e)}", show_alert=True)
